from pathlib import Path
import pymupdf
from docx import Document as DocxDocument
from pptx import Presentation
from langchain_core.documents import Document

def load_pdf(file_path:str)->list[Document]:
    documents=[]
    pdf=pymupdf.open(file_path)
    for page_number, page in enumerate(pdf,start=1):
        text=page.get_text().strip()
        if text:
            documents.append(
                Document(
                    page_content=text,
                    metadata={
                        "source":file_path,
                        "file_type":"pdf",
                        "page":page_number,
                    },
                )
            )
    pdf.close()
    return documents

def load_docx(file_path:str)->list[Document]:
    documents=[]
    doc=DocxDocument(file_path)
    paragraphs=[]
    for paragraph in doc.paragraphs:
        text=paragraph.text.strip()
        if text:
            paragraphs.append(text)

    full_text="\n".join(paragraphs)
    if full_text:
        documents.append(
            Document(
                page_content=full_text,
                metadata={
                    "source":file_path,
                    "file_type":"docx",
                },
            )
        )
    return documents

def load_pptx(file_path:str)->list[Document]:
    documents=[]
    presentation=Presentation(file_path)

    for slide_number,slide in enumerate(
        presentation.slides,start=1,
    ):
        texts=[]
        for shape in slide.shapes:
            if hasattr(shape,"text"):
                text=shape.text.strip()
                if text:
                    texts.append(text)

        slide_text="\n".join(texts)
        if slide_text:
            documents.append(
                Document(
                    page_content=slide_text,
                    metadata={
                        "source":file_path,
                        "file_type":"pptx",
                        "slide":slide_number,
                    },
                )
            )

    return documents


def load_document(file_path:str)->list[Document]:
    extension=Path(file_path).suffix.lower()

    if extension==".pdf":
        return load_pdf(file_path)

    if extension==".docx":
        return load_docx(file_path)

    if extension==".pptx":
        return load_pptx(file_path)

    raise ValueError(
        "unsupported file type."
        "Only PDF,DOCX and PPTX are supported"
    )